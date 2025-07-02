from pptx import Presentation
import os
import json

def extract_ppt_content(ppt_path):
    content = []
    try:
        prs = Presentation(ppt_path)
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
            if slide_text:
                content.append("\n".join(slide_text))
        return content
    except Exception as e:
        print(f"Error processing {ppt_path}: {str(e)}")
        return []

def process_all_ppts(directory):
    knowledge_base = {}
    
    # 按章节顺序处理文件
    files = sorted([f for f in os.listdir(directory) if f.endswith('.pptx')])
    
    for filename in files:
        chapter_name = filename.split('.')[0]
        ppt_path = os.path.join(directory, filename)
        content = extract_ppt_content(ppt_path)
        knowledge_base[chapter_name] = content
        print(f"Processed: {filename}")
    
    return knowledge_base

# 处理PPT目录
ppt_dir = "ppt"
knowledge_base = process_all_ppts(ppt_dir)

# 保存提取的内容
with open("knowledge_base.json", "w", encoding="utf-8") as f:
    json.dump(knowledge_base, f, ensure_ascii=False, indent=2) 